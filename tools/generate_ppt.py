#!/usr/bin/env python3
"""Generate NADS defense PPT using python-pptx"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# === Color palette ===
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE = RGBColor(0x00, 0x7A, 0xFF)
ORANGE = RGBColor(0xFF, 0x95, 0x00)
RED = RGBColor(0xFF, 0x3B, 0x30)
GREEN = RGBColor(0x34, 0xC7, 0x59)
PURPLE = RGBColor(0xAF, 0x52, 0xDE)
YELLOW = RGBColor(0xFF, 0xCC, 0x00)
GRAY = RGBColor(0x86, 0x86, 0x8B)
LIGHT_GRAY = RGBColor(0xC7, 0xC7, 0xCC)
CYAN = RGBColor(0x5A, 0xC8, 0xFA)
CARD_BG = RGBColor(0x2A, 0x2A, 0x5E)

def add_bg(slide):
    bg = slide.background; fill = bg.fill; fill.solid()
    fill.fore_color.rgb = DARK_BG

def add_title(slide, text, left=0.8, top=0.4, width=11.5, size=34):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.7))
    p = txBox.text_frame.paragraphs[0]
    p.text = text; p.font.size = Pt(size); p.font.bold = True; p.font.color.rgb = WHITE

def add_text(slide, text, left, top, width, height, size=14, color=WHITE, bold=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame; tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(size); p.font.color.rgb = color
        p.font.bold = bold; p.space_after = Pt(3)

def add_shape(slide, shape_type, left, top, width, height, color, text='', size=12, text_color=None):
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.fill.background()
    if text:
        tf = shape.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(size)
        p.font.color.rgb = text_color if text_color else WHITE; p.alignment = PP_ALIGN.CENTER

def add_arrow(slide, left, top, width, height, color=GRAY):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.fill.background()

def add_placeholder(slide, left, top, width, height, label):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x3E)
    shape.line.color.rgb = GRAY; shape.line.width = Pt(1)
    tf = shape.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = label; p.font.size = Pt(18)
    p.font.color.rgb = GRAY; p.alignment = PP_ALIGN.CENTER

# ═══════════════════════════════════════════════════
# SLIDE 1: Cover
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 0.06, BLUE)
add_text(slide, "NADS", 0.8, 1.8, 11, 1.0, size=68, color=BLUE, bold=True)
add_text(slide, "Network Attack Detection System", 0.8, 2.7, 11, 0.6, size=22, color=WHITE, bold=True)
add_text(slide, "基于混合检测架构（误用检测 + 异常检测）的入侵检测系统\n信息安全科技创新  |  2026年7月", 0.8, 3.6, 9, 1.0, size=16, color=GRAY)
add_text(slide, "8核心模块 + 2拓展引擎 | 136条检测规则 | 12种协议识别 | 7种异常检测器\n6-Tab Apple风格GUI | 准确率96% | 误报率1.72% | ~15,000行代码 | 34个开源项目参考", 0.8, 5.0, 11, 1.5, size=15, color=LIGHT_GRAY)
add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 7.44, 13.333, 0.06, BLUE)

# ═══════════════════════════════════════════════════
# SLIDE 2: TOC
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title(slide, "目录  CONTENTS")
items = [
    ("01", "系统架构", "五层管道架构 + 数据流 + 模块间接口"),
    ("02", "误用检测引擎",   "AC自动机 + 136条规则 + 三阶段匹配流水线"),
    ("03", "加密流量检测 (PDF必做1)", "JA3/JA4指纹 + 证书异常 + C2 Beacon CV检测"),
    ("04", "异常检测 + ML (PDF必做2)", "7种检测器 + Isolation Forest + 两阶段精判"),
    ("05", "攻击链+威胁情报 (PDF必做3+4)", "MITRE ATT&CK 7阶段串联 + AbuseIPDB/OTX"),
    ("06", "GUI + 性能评估", "6-Tab界面 + 攻击链可视化 + 准确率96%"),
    ("07", "混合检测验证 + 特色创新", "双引擎互印证 + 真实流量 + 6大特色方向"),
]
for i, (num, title, desc) in enumerate(items):
    y = 1.5 + i * 0.85
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.8, y, 0.7, 0.55, BLUE, num, 18)
    add_text(slide, title, 1.7, y, 5, 0.35, size=17, color=WHITE, bold=True)
    add_text(slide, desc, 1.7, y+0.32, 8, 0.3, size=12, color=GRAY)

# ═══════════════════════════════════════════════════
# SLIDE 3: Architecture
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title(slide, "系统架构 — 五层管道设计")

layers = [
    ("Layer 1: 数据采集 (B)", "PacketCapture | Scapy+Npcap | BPF+PCAP回放", BLUE),
    ("Layer 2: 协议解析 (B)", "ProtocolParser + TCPReassembler + TLSDetector + C2BeaconDetector | 12种协议 + TCP严格重组", CYAN),
    ("Layer 3: 检测引擎 (A+C)", "MisuseDetector(AC+136规则) + AnomalyDetector(7检测器) + MLAnomalyDetector + AttackChainAnalyzer", PURPLE),
    ("Layer 4: 告警管理 (D)", "AlertManager | 去重+分级+统计+JSON | ThreatIntel字段", ORANGE),
    ("Layer 5: 用户界面 (D)", "IDSMainWindow(PyQt5) | Apple风格6-Tab | Dashboard/Alerts/Stats/Signatures/AttackChain/Log", GREEN),
]
for i, (name, desc, color) in enumerate(layers):
    y = 1.3 + i * 1.15
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.8, y, 3.0, 0.95, color, name, 11)
    add_text(slide, desc, 4.0, y+0.1, 5.5, 0.8, size=12, color=LIGHT_GRAY)
    if i < 4:
        add_arrow(slide, 2.0, y+0.97, 0.25, 0.15, color)

stats = ["8个核心模块", "2个拓展引擎(ML+攻击链)", "136条检测规则", "12种协议识别", "7种异常检测器", "6-Tab Apple GUI", "~15,000行代码", "34个开源项目参考"]
for i, s in enumerate(stats):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 9.5, 1.5+i*0.55, 3.2, 0.45, CARD_BG, s, 13, WHITE)

# ═══════════════════════════════════════════════════
# SLIDE 4: Data Flow
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title(slide, "数据流 — 一个数据包的七步处理")

flow = [("PCAP File", BLUE), ("ProtocolParser", CYAN), ("MisuseDetector", PURPLE), ("AnomalyDetector", ORANGE), ("AlertManager", GREEN), ("GUI Dashboard", RED)]
for i, (label, color) in enumerate(flow):
    x = 0.5 + i * 2.1
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 2.0, 1.8, 0.9, color, label, 14, WHITE)
    if i < 5: add_arrow(slide, x+1.85, 2.25, 0.2, 0.4, color)

details = (
    "Step 1: ProtocolParser.parse()      IP/TCP/HTTP解析, 12种协议识别, HTTP深度解析(method/URI/Host/UA)\n"
    "Step 2: MisuseDetector.match_packet()  AC自动机(115模式)三阶段: 字符串扫描->正则精匹配->阈值判定\n"
    "Step 3: TLSDetector.analyze_client_hello()  JA3/JA4指纹提取 + 恶意指纹库匹配 + 证书异常检测\n"
    "Step 4: TCPReassembler.feed()          严格seq排序, 重传去重, 乱序缓存, LRU淘汰\n"
    "Step 5: AnomalyDetector.update()       15维HostStats更新, 每5秒周期检查7种检测器\n"
    "Step 6: MLAnomalyDetector.predict()     8维特征向量, IsolationForest无监督异常检测, 两阶段精判\n"
    "Step 7: AlertManager.submit()          去重->分级(critical/high/medium/low)->回调推送GUI->JSON导出"
)
add_text(slide, details, 0.8, 3.5, 11.5, 3.5, size=14, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════
# SLIDE 5: Misuse Detection
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title(slide, "误用检测引擎 — Aho-Corasick 多模式匹配 (成员A)")

add_text(slide, "AC自动机 = Trie树 + Failure Link", 0.8, 1.5, 6, 0.4, size=18, color=BLUE, bold=True)
ac_text = (
    "将所有136条规则的固定字符串建为一棵前缀树(Trie):\n"
    "  Goto函数:   沿树逐个字符匹配, O(1) per char\n"
    "  Failure函数: 失配时跳转到最长可匹配后缀状态 (KMP思想推广)\n"
    "  Output函数:  继承failure link的输出, 避免遗漏嵌入匹配\n\n"
    "复杂度: O(n+m) -- n=载荷长度, m=模式总长度\n"
    "        与规则数量无关! 一次遍历找出所有命中\n\n"
    "Fallback: pyahocorasick未安装时, 自动用re.escape()将\n"
    "          简单字符串转为正则匹配, 零规则丢弃"
)
add_text(slide, ac_text, 0.8, 2.0, 6, 3.2, size=13, color=LIGHT_GRAY)

add_text(slide, "三阶段匹配流水线", 7.5, 1.5, 5, 0.4, size=18, color=BLUE, bold=True)
stages = [
    ("Stage 1: AC自动机扫描", "115个固定字符串\nO(n)一次遍历, 95%+流量在此过滤", BLUE),
    ("Stage 2: 正则精匹配", "192个预编译正则\n仅对Stage1命中候补执行", PURPLE),
    ("Stage 3: 阈值型规则判定", "端口扫描/暴力破解/DDoS\n检查累积计数是否达阈值", ORANGE),
]
for i, (title, desc, color) in enumerate(stages):
    y = 2.0 + i * 1.3
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 7.5, y, 2.2, 0.45, color, title[:20], 10, WHITE)
    add_text(slide, desc, 9.9, y, 3, 1.1, size=12, color=LIGHT_GRAY)

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.8, 5.6, 11.5, 1.6, CARD_BG)
perf = (
    "性能: 136规则 | AC 115模式 + 正则 192模式 | 正常包 468us/pkt (2138 pkts/s) | 攻击包 455us/pkt\n"
    "准确率: 96% (22/23) | SQLi 80%/XSS 100%/Web 100%/BruteForce 100%/Backdoor 100%/Scan 100%/WebShell 100%\n"
    "误报率: 0% (自研规则在正常流量中零误报, 仅1条Suricata导入规则在930包中产生16次误报)\n"
    "特色: 自研Suricata规则导入器 -> 兼容Emerging Threats Open社区30,000+规则"
)
add_text(slide, perf, 1.2, 5.8, 11, 1.4, size=13, color=WHITE)

# ═══════════════════════════════════════════════════
# SLIDE 6: TLS Detection
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title(slide, "加密流量检测 — JA3/JA4指纹 + C2 Beacon CV检测 (成员B)")

add_text(slide, "JA3 = MD5(Version + CipherSuites + Extensions + Curves + Formats)", 0.8, 1.3, 12, 0.4, size=18, color=BLUE, bold=True)
ja3_text = (
    "TLS ClientHello 握手消息中的5个明文字段  ->  MD5  ->  32字符指纹\n\n"
    "不同TLS库(OpenSSL/Golang/Mozilla NSS/Cobalt Strike)有不同的参数组合 -> 如同浏览器UA指纹\n\n"
    "恶意指纹库: Trickbot(6734f374) | Emotet(51c64c77) | CobaltStrike(72a589da) | Meterpreter(e35df3e0) | C2 Demo(6eeb1f82)\n\n"
    "证书异常检测: 自签名(Issuer==Subject) | 过期(notAfter<now) | 弱加密套件(RC4/DES/3DES/NULL等15种) | CN不匹配"
)
add_text(slide, ja3_text, 0.8, 1.8, 6, 2.8, size=13, color=LIGHT_GRAY)

add_text(slide, "C2 Beacon 心跳检测 (CV算法)", 7.5, 1.3, 5.5, 0.4, size=18, color=ORANGE, bold=True)
cv_text = (
    "原理: 变异系数 CV = std/mean (标准差/均值)\n\n"
    "对每个(src_ip,dst_ip,dst_port)三元组:\n"
    "  记录每次连接时间戳 -> 计算时间间隔\n"
    "  CV<0.05 -> CRITICAL (几乎完美定时)\n"
    "  CV<0.15 -> HIGH (轻微抖动)\n"
    "  CV<0.30 -> MEDIUM (边界值)\n"
    "  CV>0.30 -> 正常(人类行为)\n\n"
    "参考: github.com/Mithileshan/c2-beaconing-detection\n"
    "       github.com/thousaba/beacon_hunter"
)
add_text(slide, cv_text, 7.5, 1.8, 5.5, 3.0, size=12, color=LIGHT_GRAY)

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.8, 5.2, 11.5, 1.8, CARD_BG)
bottom = (
    "三层防线: JA3/JA4指纹(握手指纹) + 证书异常(自签名/过期/弱加密) + C2 Beacon CV检测(时间规律)\n\n"
    "核心优势: 完全不需要解密HTTPS载荷 -- 仅凭明文ClientHello参数和连接统计即可发现恶意C2通信\n\n"
    "参考: github.com/salesforce/ja3 | github.com/FoxIO-LLC/ja4 | sslbl.abuse.ch | ja3er.com"
)
add_text(slide, bottom, 1.2, 5.4, 11, 1.5, size=14, color=WHITE)

# ═══════════════════════════════════════════════════
# SLIDE 7: Anomaly + ML
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title(slide, "异常检测 + ML — 7种检测器 + Isolation Forest (成员C)")

add_text(slide, "7种异常检测器", 0.8, 1.3, 5, 0.4, size=18, color=ORANGE, bold=True)
detectors = [
    ("端口扫描", "30端口SYN -> medium告警"),
    ("横向扫描", "60个不同目标IP -> high告警"),
    ("SYN Flood", "1200个SYN (>80%) -> high告警"),
    ("暴力破解", "6次登录失败 -> high告警"),
    ("高频DDoS", "50源x300包=15000包 -> critical告警"),
    ("基线偏离", "当前指标>基线5倍 -> 告警触发"),
    ("多源并发", "50源IP高频并发 -> DDoS告警"),
]
for i, (name, desc) in enumerate(detectors):
    y = 1.8 + i * 0.6
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.8, y, 2.4, 0.48, ORANGE, name, 12, WHITE)
    add_text(slide, desc, 3.4, y+0.05, 4, 0.42, size=12, color=LIGHT_GRAY)

add_text(slide, "Isolation Forest + TwoStageDetector", 7.5, 1.3, 5.5, 0.4, size=18, color=GREEN, bold=True)
ml_text = (
    "IF原理: 异常数据 '少且不同' -- 随机切几刀就能隔离\n"
    "8维特征向量: [连接数,SYN数,唯一端口,唯一IP,收发字节,包速率,登录失败]\n\n"
    "TwoStageDetector两阶段:\n"
    "  Stage 1: IF粗筛 -> 标记异常主机 (-1=异常)\n"
    "  Stage 2: 三源精判(投票制)\n"
    "    MisuseDetector有告警? -> +1票\n"
    "    BaselineProfile偏离>2sigma? -> +1票\n"
    "    ThreatIntel已知恶意? -> +1票\n"
    "    票数>=2 -> 确认告警 | =1 -> 低优先级 | =0 -> 抑制(误报)\n\n"
    "参考: github.com/sarthakghavghave (88.6%误报消减)\n"
    "      github.com/fatemak04/Anomaly-NIDS (IF+Scapy+基线)"
)
add_text(slide, ml_text, 7.5, 1.8, 5.5, 4.0, size=12, color=LIGHT_GRAY)

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.8, 6.3, 11.5, 0.9, CARD_BG)
add_text(slide, "稳定性: 10,000包/306k pps/内存+2.9MB | 动态阈值自适应: mu+ksigma, k=3.0, min_samples=10 | 7/7检测器触发验证 | 阈值调优:5档全测", 1.2, 6.45, 11, 0.7, size=13, color=WHITE)

# ═══════════════════════════════════════════════════
# SLIDE 8: Attack Chain + Threat Intel
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title(slide, "攻击链关联 (PDF必做3) + 威胁情报 + 降噪 (PDF必做4)")

add_text(slide, "MITRE ATT&CK 攻击链关联分析", 0.8, 1.4, 6.5, 0.4, size=18, color=RED, bold=True)
chain_text = (
    "AttackChainAnalyzer: 同源IP, 30min窗口, 按ATT&CK阶段串联\n\n"
    "阶段映射: scan/port_scan->reconnaissance | sql_injection/xss/web_attack->exploit\n"
    "          brute_force/webshell->initial_access | backdoor/trojan->persistence\n"
    "          c2/c2_beacon->command_control | lateral_scan->lateral_movement | dos->impact\n\n"
    "判定: >=3个不同阶段 -> 输出attack_chain告警\n"
    "验证: 4步攻击(scan->bruteforce->sqli->backdoor)正确串联为[侦察->初始入侵->利用->持久化]"
)
add_text(slide, chain_text, 0.8, 1.9, 6.5, 3.0, size=12, color=LIGHT_GRAY)

add_text(slide, "威胁情报 + 误报降噪", 7.5, 1.4, 5.5, 0.4, size=18, color=PURPLE, bold=True)
ti_text = (
    "威胁情报集成:\n"
    "  AbuseIPDB API (免费1000次/天, abuseConfidenceScore)\n"
    "  AlienVault OTX (免费, pulse_count)\n"
    "  本地黑名单(5个恶意IP+5个恶意JA3指纹)\n"
    "  enrich_alert()自动标注 -> 已知恶意IP升级severity\n\n"
    "误报自动降噪(三级体系):\n"
    "  1. 上下文过滤(A) -- 攻击特征位置识别\n"
    "     URL参数->正常 | Referer/Cookie->降级 | 白名单IP->忽略\n"
    "  2. 基线降噪(C) -- 告警指标在mu+-2sigma内->降级\n"
    "  3. 资产重要性(C) -- config.yaml四级配置\n"
    "     critical资产+告警->升级 | normal资产+告警->降级"
)
add_text(slide, ti_text, 7.5, 1.9, 5.5, 4.0, size=11, color=LIGHT_GRAY)

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.8, 5.8, 11.5, 1.4, CARD_BG)
add_text(slide, "攻击链验证: 4步攻击正确串联 | 威胁情报: enrich_alert() + enrich_alerts() | 降噪: config.yaml四级资产配置(critical/important/normal/whitelist) + mu+-2sigma基线过滤", 1.2, 6.0, 11, 1.0, size=13, color=WHITE)

# SLIDES 9-10: Screenshot placeholders
for slide_idx, (title, label) in enumerate([
    ("GUI 仪表盘 — Apple 风格 6-Tab 设计",
     "Dashboard 全页截图\n\n启动GUI -> 回放demo_hybrid_detection.pcap -> 截图\n\n应显示: 6张统计卡片 | PPS/BPS双线折线图 | Recent Alerts表格(含Source列anomaly+misuse) | Top Attack Sources"),
    ("GUI 攻击链可视化面板",
     "Attack Chain Tab 截图\n\n回放demo_hybrid_detection.pcap -> 切换到Attack Chain -> 点击Refresh\n\n应显示: 蓝色源IP节点 -> 黄色(侦察) -> 橙色(利用) -> 红色(C2/持久化)\n底部统计: X chains | Y steps | max Z phases | Level: HIGH/CRITICAL"),
]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, title)
    add_placeholder(slide, 0.5, 1.5, 12.3, 5.5, label)

# SLIDE 11: Performance Data
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title(slide, "性能数据汇总")

metrics = [
    ("96%", "七类攻击\n总准确率", BLUE),
    ("1.72%", "正常流量\n误报率", GREEN),
    ("468 us", "单包匹配\n平均延迟", PURPLE),
    ("306,756", "异常检测\nPPS吞吐", ORANGE),
    ("2,138", "误用检测\npkts/sec", RED),
    ("77", "8 PCAP\n总去重告警", YELLOW),
]
for i, (num, label, color) in enumerate(metrics):
    x = 0.5 + i * 2.1
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 1.3, 1.9, 1.6, color)
    add_text(slide, num, x+0.05, 1.35, 1.8, 0.6, size=26, color=WHITE, bold=True)
    add_text(slide, label, x+0.05, 2.0, 1.8, 0.7, size=11, color=WHITE)

detail = (
    "规则覆盖: 136条(自研78+Suricata58) | AC 115模式 + 正则 192模式 | 3步匹配: AC扫描->正则精匹->阈值判定\n"
    "协议识别: 12种协议 | 识别准确率100% (构造包) | TCP重组: 顺序/乱序/重传/重叠 4/4通过\n"
    "异常检测: 7种检测器全覆盖 | 10,000包稳定性验证 +2.9MB | 动态阈值自适应 mu+-3sigma\n"
    "ML模型: Isolation Forest (50样本训练) | 8维特征向量 | TwoStageDetector: IF粗筛->三源精判\n"
    "学术评估: CIC-IDS-2017 Wednesday: 17,244包->99告警(DoS50/扫描28/Web21) + 2攻击链\n"
    "真实流量: Kali->Windows扫描837包(双引擎互印证) | 正常浏览930包(误报率1.72%, 自研规则零误报)\n"
    "GUI: 6-Tab Apple风格(亮暗双主题) | 攻击链可视化(QGraphicsView) | 所有操作<200ms | 内存<500MB\n"
    "特色: Suricata规则导入 | C2 Beacon CV检测 | TwoStage检测 | 威胁情报 | 攻击链可视化 | 一键演示"
)
add_text(slide, detail, 0.8, 3.4, 11.5, 3.8, size=14, color=LIGHT_GRAY)

# SLIDE 12: Hybrid Detection
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title(slide, "混合检测双引擎验证 — 同一次攻击, 双引擎互印证")

# Left: demo_hybrid
add_text(slide, "demo_hybrid_detection.pcap (55包, 3秒)", 0.8, 1.3, 7, 0.4, size=18, color=BLUE, bold=True)
hybrid_text = (
    "同一来源IP (10.0.7.77) 四阶段攻击:\n\n"
    "  Phase 1: 端口扫描 -- 26个不同端口的SYN探测\n"
    "  Phase 2: SQL注入 -- 5种变体(UNION SELECT/OR 1=1/DROP TABLE/报错注入/SQLMap UA)\n"
    "  Phase 3: SSH暴力破解 -- 8次登录失败\n"
    "  Phase 4: 后门C2 -- AntSword UA后门连接\n\n"
    "检出结果: 13条去重告警\n\n"
    "  source=anomaly (1条, 橙色):\n"
    "    port_scan -- 异常检测发现26端口偏离基线(阈值20)\n"
    "    异常检测: '我不知道Nmap长什么样 -- 我只看到行为不正常'\n\n"
    "  source=misuse (12条, 蓝色):\n"
    "    SQLI-001/002/004/008/010/011 -- 误用检测精确匹配8条攻击特征\n"
    "    误用检测: '我不知道什么算正常 -- 但我认识这些攻击特征'"
)
add_text(slide, hybrid_text, 0.8, 1.8, 6.5, 5.0, size=12, color=LIGHT_GRAY)

# Right: verification
add_text(slide, "真实流量验证", 7.5, 1.3, 5.5, 0.4, size=18, color=ORANGE, bold=True)
verify_text = (
    "Kali -> Windows 真实扫描:\n"
    "  VirtualBox真实抓包, Kali使用Nmap扫描Windows\n"
    "  837个TCP包, 12种协议识别\n"
    "  双引擎告警: misuse=5 + anomaly=2\n"
    "  攻击链: 1条(侦察阶段)\n\n"
    "正常浏览流量误报率评估:\n"
    "  930个真实HTTP请求(网页浏览/API/图片/搜索)\n"
    "  误报: 16条 (1.72%)\n"
    "  全部来自1条Suricata社区规则(SUR-1000302, 'curl')\n"
    "  自研78条规则: 零误报\n\n"
    "结论:\n"
    "  1. 双引擎从不同维度同时告警, 互相印证\n"
    "  2. 异常检测发现未知行为模式\n"
    "  3. 误用检测精确定位已知攻击类型\n"
    "  4. 混合架构=互补短板, 提升整体检测能力"
)
add_text(slide, verify_text, 7.5, 1.8, 5.5, 5.0, size=12, color=LIGHT_GRAY)

# SLIDE 13: Academic
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title(slide, "学术评估 — CIC-IDS-2017 标准数据集")

add_text(slide, "CIC-IDS-2017 Wednesday DoS 子集", 0.8, 1.3, 7, 0.4, size=20, color=BLUE, bold=True)
cids = (
    "数据集: 加拿大UNB大学 CIC-IDS-2017 (学术界最常用的IDS评估基准)\n"
    "子集: Wednesday DoS/DDoS (Slowloris, Slowhttptest, Hulk, GoldenEye)\n"
    "规模: 17,244个TCP包, 36个唯一IP, 15,262 HTTP + 910 HTTPS + 762 Unknown\n\n"
    "NADS检出结果:\n"
    "  DoS/高频流量: 50条 (Slowloris不完全头 + HTTP Flood + Range攻击)\n"
    "  扫描探测:     28条 (端口扫描 + 横向扫描)\n"
    "  Web攻击:      21条 (RCE Netcat + PHP代码执行 + 文件包含)\n"
    "  总计: 99条去重告警 + 2条攻击链\n\n"
    "检测类型: dos(50), scan(28), web_attack(21)\n"
    "告警严重度: critical(37), high(45), medium(17)"
)
add_text(slide, cids, 0.8, 1.8, 7.5, 4.5, size=13, color=LIGHT_GRAY)

add_placeholder(slide, 8.5, 1.8, 4.3, 2.5, "Statistics页截图\n(回放wednesday_subset后\n饼图+柱状图)")

add_text(slide, "正常流量误报率评估", 0.8, 6.3, 6, 0.4, size=18, color=GREEN, bold=True)
add_text(slide, "normal_real_browsing.pcap: 930个真实网页浏览HTTP请求 -> 16条误报(1.72%), 全部来自1条Suricata社区规则(curl), 自研78条规则零误报", 0.8, 6.7, 12, 0.5, size=13, color=LIGHT_GRAY)

# SLIDE 14: Open Source
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title(slide, "开源项目调研 — 34个项目参考")

refs = [
    ("ThreatWire",           "AC+1200规则+Suricata",   "A"),
    ("NetGuard IDS",         "10-Tab GUI+ML+威胁地图",   "A,D"),
    ("Slips",                "行为ML+威胁情报+Zeek",    "C"),
    ("sarthakghavghave",     "IF+RF两阶段88.6%降噪",    "C"),
    ("NetPortMon",           "PyQt5实时图表",            "D"),
    ("fatemak04/Anomaly",    "实时IF+Scapy+基线",        "C"),
    ("salesforce/ja3",       "JA3标准实现",              "B"),
    ("FoxIO/ja4",            "JA4新一代指纹",            "B"),
    ("c2-beacon-detection",  "CV变异系数C2检测",         "B"),
    ("beacon_hunter",        "PCAP C2 Beacon检测",       "B"),
    ("surinort-ast",         "Suricata规则AST解析器",    "A"),
    ("SigmaForge",           "YARA<->Sigma<->Suricata", "A"),
    ("soc-intelhub",         "VT+IPDB+OTX+MITRE",       "C,D"),
    ("threat-intel-agg",     "威胁情报聚合+热力图",       "C"),
    ("Deadfall",             "BloodHound风格PCAP可视化", "D"),
    ("temporal-ids-bench",   "CIC-IDS2017标准评估",      "ALL"),
]
for i, (name, desc, member) in enumerate(refs):
    col = i // 8; row = i % 8
    x = 0.5 + col * 6.3; y = 1.5 + row * 0.7
    mc = {'A': BLUE, 'B': GREEN, 'C': ORANGE, 'D': PURPLE}.get(member.split(',')[0], GRAY)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 0.35, 0.4, mc, member[:1], 11, WHITE)
    add_text(slide, name, x+0.45, y, 3, 0.25, size=12, color=WHITE, bold=True)
    add_text(slide, desc, x+0.45, y+0.22, 3, 0.2, size=10, color=GRAY)

add_text(slide, "+18 more projects (详见 docs/特色拓展方向与开源参考.md)", 0.5, 7.0, 12, 0.3, size=11, color=GRAY)

# SLIDE 15: 6 Innovations
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title(slide, "六大特色创新")

innovations = [
    ("1 Suricata规则导入器 (A)", BLUE,
     "自研规则转换工具, 兼容Emerging Threats Open社区\n规则数: 78 -> 136 -> 可扩展至30,000+\n基于 surinort-ast + SigmaForge"),
    ("2 C2 Beaconing心跳检测 (B)", GREEN,
     "CV变异系数检测C2通信规律, 不解密TLS即可发现\nCV<0.05->critical, CV<0.15->high\n参考: c2-beaconing-detection + beacon_hunter"),
    ("3 IF+规则两阶段检测 (C)", ORANGE,
     "IF粗筛 -> 三源精判(检测+基线+情报投票)\n参考 88.6%误报消减方案\nTwoStageDetector: IF + RF完整实现"),
    ("4 威胁情报集成 (C+D)", PURPLE,
     "AbuseIPDB + AlienVault OTX + 本地黑名单\n5个恶意IP + 5个恶意JA3指纹\nenrich_alert()自动标注 -> 升级已知恶意IP告警"),
    ("5 攻击链可视化面板 (D)", RED,
     "QGraphicsView 节点-边有向图, MITRE ATT&CK颜色标记\n黄=侦察 -> 橙=利用 -> 红=C2/持久化\n同一IP多步攻击自动串联 -> 威胁等级判定"),
    ("6 一键全自动演示模式 (D)", YELLOW,
     "Demo模式: 9类攻击载荷(47个) + 50%正常流量混合\nTrafficGenerator实时生成, 无需PCAP文件\n大字幕+Toast通知+统计摘要弹窗"),
]
for i, (title, color, desc) in enumerate(innovations):
    row = i // 2; col = i % 2
    x = 0.5 + col * 6.3; y = 1.2 + row * 3.05
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 5.8, 2.8, color)
    add_text(slide, title, x+0.3, y+0.1, 5, 0.35, size=16, color=WHITE, bold=True)
    add_text(slide, desc, x+0.3, y+0.55, 5, 2.1, size=13, color=WHITE)

# SLIDE 16: Thank You
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 0.06, BLUE)
add_text(slide, "Thank You", 0.8, 1.8, 11, 1.0, size=68, color=BLUE, bold=True)
add_text(slide, "Q & A", 0.8, 2.8, 11, 0.6, size=34, color=WHITE)
add_text(slide, "NADS -- 基于混合检测架构的常见网络攻击检测系统", 0.8, 4.0, 11, 0.5, size=20, color=GRAY)
summary = (
    "误用检测 96%准确率 + 异常检测 7种检测器 + ML无监督学习 + 加密流量三层防线\n"
    "攻击链关联 + 威胁情报 + Apple风格6-Tab GUI + 真实流量双引擎验证\n"
    "~15,000行代码 | 34个开源项目参考 | 6大特色创新 | CIC-IDS-2017学术评估"
)
add_text(slide, summary, 0.8, 4.8, 11, 1.5, size=16, color=LIGHT_GRAY)
add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 7.44, 13.333, 0.06, BLUE)

# Save
out = 'docs/NADS_答辩PPT.pptx'
prs.save(out)
print(f"Saved: {out} ({len(prs.slides)} slides)")
